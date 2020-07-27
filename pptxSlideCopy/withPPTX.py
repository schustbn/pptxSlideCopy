from pptx import Presentation
import opcdiag

master = '/Users/bjoernschuster/Downloads/PPTmerge/master.pptx'
part1 = '/Users/bjoernschuster/Downloads/PPTmerge/simpel.pptx'

master_prs = Presentation(master)
part1_prs = Presentation(part1)

print(len(part1_prs.slides))

print(part1_prs.slide_master.slide_layouts[0].slide_master)

for slide in part1_prs.slides:
    print(slide.slide_id)
    print(slide.slide_layout.slide_master.name)



